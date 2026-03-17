#!/usr/bin/env python3
"""Generate PowerPoint slides for both AI talks."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Design constants ──────────────────────────────────────────────
BG_COLOR = RGBColor(0x1A, 0x1E, 0x2E)
BG_LIGHTER = RGBColor(0x2A, 0x2E, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
CYAN = RGBColor(0x00, 0xD4, 0xF5)
DIM_WHITE = RGBColor(0xDD, 0xDD, 0xEE)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_NAME = "Calibri"


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color=BG_LIGHTER, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=24, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_para(tf, text, size=20, color=WHITE, bold=False, space_before=Pt(6), bullet=False, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    if bullet:
        p.level = 1
    return p


def add_bullet_slide(prs, title_text, bullets, subtitle=None, accent=ORANGE, title_size=36, bullet_size=22):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.6), Inches(0.08), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    
    # Title
    txb = tb(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.9))
    set_text(txb.text_frame, title_text, size=title_size, bold=True, color=WHITE)
    
    # Subtitle
    y_start = Inches(1.6)
    if subtitle:
        txb2 = tb(slide, Inches(1.1), Inches(1.4), Inches(10), Inches(0.5))
        set_text(txb2.text_frame, subtitle, size=18, color=LIGHT_GRAY)
        y_start = Inches(2.0)
    
    # Bullets
    txb3 = tb(slide, Inches(1.1), y_start, Inches(11), Inches(5.0))
    tf = txb3.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(10)
        
        if isinstance(bullet, tuple):
            # (title, description)
            p.text = f"▸ {bullet[0]}"
            p.font.size = Pt(bullet_size)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = FONT_NAME
            if bullet[1]:
                p2 = tf.add_paragraph()
                p2.text = f"    {bullet[1]}"
                p2.font.size = Pt(bullet_size - 4)
                p2.font.color.rgb = LIGHT_GRAY
                p2.font.name = FONT_NAME
                p2.space_before = Pt(2)
        else:
            p.text = f"▸ {bullet}"
            p.font.size = Pt(bullet_size)
            p.font.color.rgb = WHITE
            p.font.bold = False
            p.font.name = FONT_NAME
    
    return slide


def add_table_slide(prs, title_text, headers, rows, accent=CYAN, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    # Title
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.4), Inches(0.08), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    
    txb = tb(slide, Inches(0.9), Inches(0.3), Inches(11), Inches(0.7))
    set_text(txb.text_frame, title_text, size=32, bold=True, color=WHITE)
    
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_w = Inches(11.8)
    table_h = Inches(5.5)
    
    tbl_shape = slide.shapes.add_table(num_rows, num_cols, Inches(0.75), Inches(1.2), table_w, table_h)
    table = tbl_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    
    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = FONT_NAME
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x25, 0x2A, 0x3A)
    
    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = DIM_WHITE
                paragraph.font.name = FONT_NAME
                if j == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = ORANGE
            cell.fill.solid()
            if i % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0x1E, 0x22, 0x32)
            else:
                cell.fill.fore_color.rgb = BG_COLOR
    
    # Remove table borders
    for i in range(num_rows):
        for j in range(num_cols):
            cell = table.cell(i, j)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
    
    return slide


# ═══════════════════════════════════════════════════════════════════
#  TALK 1: AI in 2026 — Overview
# ═══════════════════════════════════════════════════════════════════

def create_talk1():
    prs = new_prs()
    
    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    # Decorative accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.4), Inches(1.5), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()
    
    txb = tb(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2))
    set_text(txb.text_frame, "AI in 2026", size=54, bold=True, color=WHITE)
    add_para(txb.text_frame, "What You Need to Know", size=40, color=CYAN, bold=False, space_before=Pt(4))
    
    txb2 = tb(slide, Inches(1.5), Inches(4.5), Inches(8), Inches(0.5))
    set_text(txb2.text_frame, "Enriched Solutions  •  Ali Sherazi", size=22, color=LIGHT_GRAY)
    
    # Tagline
    txb3 = tb(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.5))
    set_text(txb3.text_frame, '"You cannot change the wind, but working together, we can adjust the sails"', size=14, color=RGBColor(0x88, 0x88, 0x99), font_name="Calibri")
    txb3.text_frame.paragraphs[0].font.italic = True
    
    # ── Slide 2: The Big Picture ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.6), Inches(0.08), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    
    txb = tb(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.9))
    set_text(txb.text_frame, "The Big Picture", size=38, bold=True, color=WHITE)
    
    txb2 = tb(slide, Inches(1.1), Inches(1.5), Inches(10), Inches(0.6))
    set_text(txb2.text_frame, "Models are the engine. Tools are the car.", size=24, color=ORANGE, bold=True)
    
    # Three concept boxes
    concepts = [
        ("🏢  Companies", "OpenAI, Google, Anthropic,\nMeta, xAI, Perplexity", "Build the AI"),
        ("🧠  Models", "GPT-5, Gemini, Claude,\nLlama 4, Grok 3", "The intelligence"),
        ("🔧  Tools", "ChatGPT, Cursor, Copilot,\nclaude.ai, Perplexity", "How you use it"),
    ]
    
    for i, (title, items, desc) in enumerate(concepts):
        x = Inches(1.1 + i * 3.8)
        box = add_shape_bg(slide, x, Inches(2.5), Inches(3.4), Inches(3.8), BG_LIGHTER, border=RGBColor(0x3A, 0x3E, 0x4E))
        
        txb = tb(slide, x + Inches(0.3), Inches(2.7), Inches(2.8), Inches(0.6))
        set_text(txb.text_frame, title, size=22, bold=True, color=CYAN)
        
        txb = tb(slide, x + Inches(0.3), Inches(3.4), Inches(2.8), Inches(1.5))
        set_text(txb.text_frame, items, size=18, color=DIM_WHITE)
        
        txb = tb(slide, x + Inches(0.3), Inches(5.2), Inches(2.8), Inches(0.5))
        set_text(txb.text_frame, desc, size=16, color=ORANGE, bold=True)
    
    # Arrow indicators
    for i in range(2):
        x = Inches(1.1 + (i + 1) * 3.8 - 0.5)
        txb = tb(slide, x, Inches(3.8), Inches(0.5), Inches(0.5))
        set_text(txb.text_frame, "→", size=30, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    
    txb_note = tb(slide, Inches(1.1), Inches(6.6), Inches(10), Inches(0.5))
    set_text(txb_note.text_frame, "Many tools let you swap models — like choosing which engine goes in your car.", size=16, color=LIGHT_GRAY)
    txb_note.text_frame.paragraphs[0].font.italic = True
    
    # ── Slide 3: AI Landscape Table ──
    add_table_slide(prs, "The AI Landscape: Who's Who",
        ["Company", "Model", "Tool", "Known For"],
        [
            ["OpenAI", "GPT-4o, GPT-5, o3", "ChatGPT", "The household name. Best all-rounder"],
            ["Google", "Gemini 2.5, Gemini 3", "Gemini app", "Baked into Gmail, Docs, Sheets"],
            ["Anthropic", "Claude Opus 4, Sonnet 4", "claude.ai", "Long docs, analysis, careful reasoning"],
            ["Meta", "Llama 4 (open-source)", "Many platforms", "Free and open — anyone can build with it"],
            ["xAI", "Grok 3", "X (Twitter)", "Real-time info, social media integrated"],
            ["Perplexity", "Multiple models", "perplexity.ai", "Search with citations — answers, not links"],
        ],
        col_widths=[1.8, 2.8, 2.2, 5.0]
    )
    
    # ── Slide 4: Coding & Dev Tools ──
    add_table_slide(prs, "Coding & Development Tools",
        ["Tool", "What It Is", "How It Works", "Why It Matters"],
        [
            ["GitHub Copilot", "AI pair programmer", "Suggests code as you type", "Most widely adopted — millions use it"],
            ["Cursor", "AI-native code editor", "Describe what you want, it writes it", "Rethought how code gets written"],
            ["Claude Code", "Anthropic's coding agent", "Reads codebase, plans & executes changes", "Handles entire features autonomously"],
            ["Codex", "OpenAI's coding agent", "Takes a task, works step-by-step", "Cloud-based, spins up own environment"],
            ["Windsurf", "AI-powered editor", "Deep AI integration for writing code", "Strong free tier, fast iteration"],
        ],
        accent=ORANGE,
        col_widths=[2.2, 2.5, 3.5, 3.6]
    )
    
    # ── Slide 5: Non-Technical AI Tools ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.6), Inches(0.08), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    
    txb = tb(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.9))
    set_text(txb.text_frame, "AI Tools for Everyone", size=36, bold=True, color=WHITE)
    
    txb_sub = tb(slide, Inches(1.1), Inches(1.3), Inches(10), Inches(0.4))
    set_text(txb_sub.text_frame, "No coding required — these tools work for any business professional", size=18, color=LIGHT_GRAY)
    
    tools = [
        ("Claude Co-work", "Collaborative AI workspace. Draft docs,\nspreadsheets, artifacts together.", CYAN),
        ("OpenClaw", "Open-source AI assistant. Connects to\nyour email, calendar, files. Runs locally.", ORANGE),
        ("Notion AI", "AI built into project management.\nAutomate notes, tasks, and wikis.", CYAN),
        ("Gamma.app", "AI presentation generation.\nText in → designed slides out.", ORANGE),
        ("Perplexity", "Research with citations.\nLike Google, but gives you answers.", CYAN),
    ]
    
    for i, (name, desc, accent_col) in enumerate(tools):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.0)
        y = Inches(2.0 + row * 2.6)
        
        box = add_shape_bg(slide, x, y, Inches(3.6), Inches(2.2), BG_LIGHTER, border=accent_col)
        
        txb = tb(slide, x + Inches(0.25), y + Inches(0.25), Inches(3.1), Inches(0.5))
        set_text(txb.text_frame, name, size=22, bold=True, color=accent_col)
        
        txb = tb(slide, x + Inches(0.25), y + Inches(0.85), Inches(3.1), Inches(1.2))
        set_text(txb.text_frame, desc, size=16, color=DIM_WHITE)
    
    # ── Slide 6: 5 Business Use Cases Overview ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.6), Inches(0.08), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()
    
    txb = tb(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.9))
    set_text(txb.text_frame, "5 Things You're About to See Live", size=36, bold=True, color=WHITE)
    
    txb_sub = tb(slide, Inches(1.1), Inches(1.3), Inches(10), Inches(0.4))
    set_text(txb_sub.text_frame, "No pre-prepared demos. No tricks. Just AI tools and real business problems.", size=18, color=LIGHT_GRAY)
    
    use_cases = [
        ("📊", "Financial Analysis", "Messy Excel → Board-ready report", "Claude"),
        ("📝", "Contract Drafting", "One brief → Full legal agreement", "Claude"),
        ("🎨", "Brand Identity", "Text description → Logo & brand guide", "Claude"),
        ("🔍", "Market Research", "A question → Sourced research report", "Perplexity + Claude"),
        ("🚀", "Pitch Deck", "Napkin idea → Designed slide deck", "Claude + Gamma"),
    ]
    
    for i, (emoji, title, desc, tool) in enumerate(use_cases):
        y = Inches(1.9 + i * 1.05)
        box = add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(0.9), BG_LIGHTER)
        
        txb = tb(slide, Inches(1.0), y + Inches(0.1), Inches(0.5), Inches(0.6))
        set_text(txb.text_frame, emoji, size=24, alignment=PP_ALIGN.CENTER)
        
        txb = tb(slide, Inches(1.6), y + Inches(0.1), Inches(2.5), Inches(0.6))
        set_text(txb.text_frame, title, size=22, bold=True, color=ORANGE)
        
        txb = tb(slide, Inches(4.3), y + Inches(0.1), Inches(5.0), Inches(0.6))
        set_text(txb.text_frame, desc, size=18, color=DIM_WHITE)
        
        txb = tb(slide, Inches(9.8), y + Inches(0.1), Inches(2.2), Inches(0.6))
        set_text(txb.text_frame, tool, size=16, color=CYAN, alignment=PP_ALIGN.RIGHT)
    
    # ── Slides 7-11: Individual Use Cases ──
    use_case_details = [
        {
            "title": "📊  Financial Analysis",
            "what": "Turn a messy Excel spreadsheet into a clean executive summary with charts and risk analysis",
            "tool": "Claude (claude.ai)",
            "see": [
                "Messy spreadsheet → cleaned & standardised data",
                "Key ratios: gross margin, EBITDA, burn rate, runway",
                "Auto-generated charts and risk matrix",
                "Board-ready summary in under 5 minutes",
            ],
            "wow": "This normally takes a financial analyst 2–3 hours.",
        },
        {
            "title": "📝  Contract Drafting",
            "what": "Generate a full professional services agreement from a one-paragraph brief, then refine clauses in real-time",
            "tool": "Claude (claude.ai)",
            "see": [
                "One-paragraph brief → multi-page legal contract",
                "12 properly structured sections with cross-references",
                "Live clause refinement (IP, termination, liability)",
                "Plain-English simplification on demand",
            ],
            "wow": "A solicitor charges £500–£2,000 for a first draft. AI gets you 80% there in 2 minutes.",
        },
        {
            "title": "🎨  Brand Identity",
            "what": "Create a logo, colour palette, brand guidelines and business card mockup from a text description",
            "tool": "Claude (claude.ai) + DALL·E",
            "see": [
                "Logo generation from a text brief",
                "Iterative refinement (more minimal, different style)",
                "Full brand colour palette with hex codes",
                "Typography recommendations & business card mockup",
            ],
            "wow": "A branding agency charges £5,000–£15,000. Get a solid starting point in minutes.",
        },
        {
            "title": "🔍  Market Research",
            "what": "Produce a structured competitive analysis with real, cited sources you can verify",
            "tool": "Perplexity Pro + Claude",
            "see": [
                "Market size, key players, trends, opportunities",
                "Every claim backed by clickable sources",
                "SWOT analysis for a specific new entrant",
                "TAM/SAM/SOM market sizing with reasoning",
            ],
            "wow": "A week-long consulting engagement compressed into minutes.",
        },
        {
            "title": "🚀  Investor Pitch Deck",
            "what": "Turn a one-line business idea into a complete 12-slide investor pitch deck with designed slides",
            "tool": "Claude (claude.ai) + Gamma.app",
            "see": [
                "Full deck: problem, solution, market, team, financials",
                "Emotional storytelling on the problem slide",
                "3-year financial projections with assumptions",
                "Designed slides via Gamma in one click",
            ],
            "wow": "Founders spend weeks on pitch decks. Get a solid first draft in minutes.",
        },
    ]
    
    for uc in use_case_details:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        
        # Title
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.6), Inches(0.08), Inches(0.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ORANGE
        bar.line.fill.background()
        
        txb = tb(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.9))
        set_text(txb.text_frame, uc["title"], size=36, bold=True, color=WHITE)
        
        # What we'll do
        txb = tb(slide, Inches(1.1), Inches(1.5), Inches(10), Inches(0.6))
        set_text(txb.text_frame, uc["what"], size=20, color=LIGHT_GRAY)
        
        # Tool badge
        badge = add_shape_bg(slide, Inches(1.1), Inches(2.2), Inches(4.5), Inches(0.55), RGBColor(0x20, 0x25, 0x35), border=CYAN)
        txb = tb(slide, Inches(1.3), Inches(2.25), Inches(4.0), Inches(0.45))
        set_text(txb.text_frame, f"Tool:  {uc['tool']}", size=18, color=CYAN, bold=True)
        
        # What you'll see
        txb = tb(slide, Inches(1.1), Inches(3.1), Inches(10), Inches(0.5))
        set_text(txb.text_frame, "What you'll see:", size=20, color=ORANGE, bold=True)
        
        txb = tb(slide, Inches(1.1), Inches(3.7), Inches(10), Inches(2.5))
        tf = txb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(uc["see"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.space_before = Pt(12)
            p.text = f"▸  {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = DIM_WHITE
            p.font.name = FONT_NAME
        
        # Wow factor
        wow_box = add_shape_bg(slide, Inches(1.1), Inches(6.0), Inches(11.0), Inches(0.7), RGBColor(0x25, 0x1A, 0x10), border=ORANGE)
        txb = tb(slide, Inches(1.4), Inches(6.05), Inches(10.4), Inches(0.6))
        set_text(txb.text_frame, f"💡  {uc['wow']}", size=18, color=ORANGE, bold=True)
    
    # ── Slide 12: Let's Try It ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    txb = tb(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.2))
    set_text(txb.text_frame, "Let's Try It", size=54, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()
    
    txb2 = tb(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1.0))
    set_text(txb2.text_frame, "5 live demos. No safety net.", size=28, color=CYAN, alignment=PP_ALIGN.CENTER)
    add_para(txb2.text_frame, "Everything you're about to see uses tools available to you today — most of them free.", 
             size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(16))
    
    return prs


# ═══════════════════════════════════════════════════════════════════
#  TALK 2: Workshop — Hands-On
# ═══════════════════════════════════════════════════════════════════

def add_prompt_slide(prs, prompt_text, title="The Prompt"):
    """Special slide for prompts — prominent text in a framed box."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    # Title
    txb = tb(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.6))
    set_text(txb.text_frame, title, size=24, bold=True, color=CYAN)
    
    # Prompt box
    box = add_shape_bg(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.8), BG_LIGHTER, border=RGBColor(0x3A, 0x3E, 0x4E))
    
    # Truncate if very long — keep the essential parts
    lines = prompt_text.strip().split('\n')
    # Show full prompt but cap at reasonable size
    display_text = prompt_text.strip()
    
    txb = tb(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.4))
    tf = txb.text_frame
    tf.word_wrap = True
    
    # Determine font size based on text length
    if len(display_text) > 800:
        font_size = 11
    elif len(display_text) > 500:
        font_size = 13
    elif len(display_text) > 300:
        font_size = 15
    else:
        font_size = 17
    
    first_line = True
    for line in display_text.split('\n'):
        if first_line:
            p = tf.paragraphs[0]
            first_line = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(2)
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xE0, 0xE8, 0xF0)
        p.font.name = "Consolas"
    
    return slide


def create_talk2():
    prs = new_prs()
    
    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.4), Inches(1.5), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()
    
    txb = tb(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2))
    set_text(txb.text_frame, "AI Workshop", size=54, bold=True, color=WHITE)
    add_para(txb.text_frame, "Hands-On with Real Business Tasks", size=36, color=ORANGE, bold=False, space_before=Pt(4))
    
    txb2 = tb(slide, Inches(1.5), Inches(4.5), Inches(8), Inches(0.5))
    set_text(txb2.text_frame, "Enriched Solutions  •  Ali Sherazi", size=22, color=LIGHT_GRAY)
    
    # ── Slide 2: Setup ──
    add_bullet_slide(prs, "Before We Start", [
        ("Claude (claude.ai)", "Our primary tool — logged in with a fresh chat"),
        ("Perplexity Pro", "For research with cited sources"),
        ("Gamma.app", "To turn text into designed slides"),
        ("A sense of adventure", "Everything is live — things might break!"),
    ], subtitle="Tools & accounts you'll need to follow along", accent=CYAN)
    
    # ═══ USE CASE 1: Financial Analysis ═══
    
    # Slide A: Scenario
    add_bullet_slide(prs, "Use Case 1: Financial Analysis", [
        "Company: GreenPulse Ltd — early-stage clean energy startup",
        "Problem: 12 months of messy P&L data — mixed currencies, missing values, bad formatting",
        "Goal: Clean it, calculate key ratios, produce a board-ready executive summary",
        "Time to do this manually: 2–3 hours with a financial analyst",
    ], subtitle="📊  Turn messy financials into an executive summary", accent=ORANGE)
    
    # Slide B: Prompt
    prompt1 = """I've uploaded a 12-month P&L spreadsheet for GreenPulse Ltd, an early-stage clean energy startup. The data is messy — inconsistent date formats, some values in GBP and others in EUR (use 1 EUR = 0.86 GBP), missing values, and inconsistent formatting.

Please:
1. Clean the data — standardise to GBP, fix formatting, flag missing values
2. Calculate: Gross margin %, EBITDA margin %, monthly burn rate, revenue growth rate, runway estimate (£420K cash)
3. Identify the 3 most important trends
4. Produce a clean summary table (12 months + annual totals)
5. Write a 200-word board-ready executive summary

Be specific with numbers — no vague statements."""
    add_prompt_slide(prs, prompt1, "Use Case 1 — The Prompt")
    
    # Slide C: Takeaways
    add_bullet_slide(prs, "Use Case 1: Key Takeaways", [
        ("Specificity is everything", "Detailed prompt → analyst-quality output"),
        ("Iterate, don't expect perfection", "Follow up: charts, risk matrix, board summary"),
        ("Context matters", "We told it the company type, exchange rate, exact metrics"),
        ("The 80% rule", "AI gets you 80% there — your expertise finishes it"),
    ], subtitle="📊  What to notice in the output", accent=ORANGE)
    
    # ═══ USE CASE 2: Contract Drafting ═══
    
    add_bullet_slide(prs, "Use Case 2: Contract Drafting", [
        "Scenario: Draft a professional services agreement for a consulting engagement",
        "Client: NovaTech Solutions (£45M manufacturing company, digital transformation)",
        "Provider: Meridian Consulting — 6 months, £75,000 + VAT",
        "Goal: Full 12-section legal contract from a one-paragraph brief",
    ], subtitle="📝  From brief to legal agreement in minutes", accent=CYAN)
    
    prompt2 = """Draft a Professional Services Agreement:

Provider: Meridian Consulting Ltd, 45 King Street, London EC2V 8QA
Client: NovaTech Solutions Ltd, 12 Innovation Park, Manchester M1 3BB
Scope: Digital transformation advisory — assessment, roadmap, vendor selection, implementation oversight
Duration: 6 months from 1 April 2026, extendable by mutual agreement
Fee: £75,000, 6 monthly instalments of £12,500 + VAT, due within 14 days
Key Personnel: Lead Consultant (10+ yrs) + Senior Analyst, no substitution without approval
Terms: Remote-first, monthly on-site Manchester visits, 30 days termination notice

Include: Definitions, Scope, Duration, Fees (+ late payment), Personnel, Client Obligations, Confidentiality, IP, Limitation of Liability, Termination, Dispute Resolution, General Provisions.

Use formal legal English. Comprehensive but not verbose. Ready for solicitor review."""
    add_prompt_slide(prs, prompt2, "Use Case 2 — The Prompt")
    
    add_bullet_slide(prs, "Use Case 2: Key Takeaways", [
        ("Structure your request", "Specifying 12 sections prevents missed clauses"),
        ("Iterate for negotiation", "Shift contract balance with a single follow-up prompt"),
        ("Plain English on demand", "Simplify legal jargon while keeping enforceability"),
        ("Still need a lawyer", "AI drafts, humans review — this saves days, not replaces expertise"),
    ], subtitle="📝  What to notice in the output", accent=CYAN)
    
    # ═══ USE CASE 3: Brand Identity ═══
    
    add_bullet_slide(prs, "Use Case 3: Brand Identity", [
        "Company: Aether — AI-powered air quality monitoring for smart cities",
        "Audience: City councils, property developers, health-conscious professionals",
        "Brand feel: Clean, scientific, trustworthy, 'calm technology'",
        "Goal: Logo, colour palette, typography, business card mockup",
    ], subtitle="🎨  Create a complete brand identity from a text description", accent=ORANGE)
    
    prompt3 = """Build a brand identity for "Aether" — an AI-powered air quality monitoring platform for smart cities. Target: city councils, property developers, health-conscious urban professionals.

Brand feel: clean, scientific, trustworthy, slightly futuristic — "calm technology."

Generate a logo:
- Modern logomark (icon + wordmark)
- Icon: abstractly represents air/atmosphere — flowing lines, gradients, stylised "A"
- Palette: deep teal/navy primary, fresh green or sky blue accent
- Clean sans-serif typeface
- Works on white and dark backgrounds, at small sizes

Style: minimal, geometric, professional. Think Stripe, Linear, Notion.
NOT cartoonish, not complex."""
    add_prompt_slide(prs, prompt3, "Use Case 3 — The Prompt")
    
    add_bullet_slide(prs, "Use Case 3: Key Takeaways", [
        ("Reference brands set the bar", "Naming Stripe/Linear/Notion guides the AI's aesthetic"),
        ("Say what you DON'T want", "'Not cartoonish, not complex' constrains the creative space"),
        ("Build incrementally", "Logo → refine → palette → typography → mockup"),
        ("A brief for your designer", "Not a finished identity — a solid starting direction"),
    ], subtitle="🎨  What to notice in the output", accent=ORANGE)
    
    # ═══ USE CASE 4: Market Research ═══
    
    add_bullet_slide(prs, "Use Case 4: Market Research", [
        "Market: UK sustainable packaging industry",
        "Tool: Perplexity Pro — real-time web search with cited sources",
        "Goal: Market size, competitors, trends, SWOT, TAM/SAM/SOM",
        "The difference: Every claim links to a real, verifiable source",
    ], subtitle="🔍  Deep research with citations, not hallucinations", accent=CYAN)
    
    prompt4 = """Comprehensive market research: UK sustainable packaging industry (2026).

1. Market Overview: Size (£), growth rate (CAGR), projected 2030 size.
   Segments: biodegradable plastics, paper/cardboard, reusable packaging, edible/dissolvable.

2. Key Drivers: Regulatory (UK Plastic Packaging Tax, EPR), consumer demand, retailer commitments (name companies), tech advances.

3. Competitive Landscape: Top 10 UK companies — focus area, revenue/funding, differentiator.

4. Challenges: Cost premiums, performance limits, supply chain, greenwashing.

5. Trends 2026-2028: Emerging materials, AI in packaging, circular economy, D2C shifts.

Cite all statistics with sources. Professional tone for a board strategy presentation."""
    add_prompt_slide(prs, prompt4, "Use Case 4 — The Prompt")
    
    add_bullet_slide(prs, "Use Case 4: Key Takeaways", [
        ("Sources make it real", "Every stat links to a verifiable article or report"),
        ("Start broad, drill deep", "Market overview → SWOT → competitive threat → TAM/SAM/SOM"),
        ("Not a McKinsey report", "But 70-80% of the way there in 90 seconds"),
        ("Verify before you present", "Click the citations — trust but verify"),
    ], subtitle="🔍  What to notice in the output", accent=CYAN)
    
    # ═══ USE CASE 5: Pitch Deck ═══
    
    add_bullet_slide(prs, "Use Case 5: Investor Pitch Deck", [
        "Startup: Aether — smart city air quality monitoring (the brand we just created!)",
        "Round: Seed, raising £4M",
        "Goal: 12-slide pitch deck — then design it with Gamma.app",
        "Two AI tools working together: Claude for content, Gamma for design",
    ], subtitle="🚀  From a napkin idea to a designed pitch deck", accent=ORANGE)
    
    prompt5 = """Create a 12-slide investor pitch deck for Aether (Seed round, raising £4M):

Aether: B2B SaaS — IoT sensors + AI analytics for real-time hyperlocal air quality monitoring. 3 pilot cities, £1.2M pre-seed, team of 4 (ex-Dyson, ex-DeepMind, ex-Capita, ex-Octopus Energy).

Slides:
1. Title — name, tagline, one-liner
2. The Problem — 36K deaths/year in UK. Emotional + urgent.
3. The Solution — What Aether does, simply.
4. How It Works — 3 simple steps.
5. Market Opportunity — size, growth, why now.
6. Traction — 3 pilots, key metrics, testimonials.
7. Business Model — pricing tiers.
8. Competitive Landscape — 2x2 matrix.
9. The Team — founders + credentials.
10. Go-to-Market — 12-month plan.
11. Financials — 3-year forecast.
12. The Ask — £4M, use of funds, milestones.

Each slide: headline, 2-4 bullet points, visual suggestion. Confident, investor-facing tone."""
    add_prompt_slide(prs, prompt5, "Use Case 5 — The Prompt")
    
    add_bullet_slide(prs, "Use Case 5: Key Takeaways", [
        ("Structure drives quality", "Specifying each slide = less editing later"),
        ("Chain your tools", "Claude writes content → Gamma designs slides → done"),
        ("Storytelling matters", "Ask for emotional framing, not just facts"),
        ("A starting point, not a final deck", "Refine with your own expertise and judgement"),
    ], subtitle="🚀  What to notice in the output", accent=ORANGE)
    
    # ── Wrap-up Slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.6), Inches(0.08), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()
    
    txb = tb(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.9))
    set_text(txb.text_frame, "What We Learned", size=38, bold=True, color=WHITE)
    
    lessons = [
        ("AI is an accelerator, not a replacement", "It gets you 70-80% there. Your expertise finishes the job."),
        ("Specificity is everything", "Vague prompt → vague answer. Detailed prompt → professional output."),
        ("Iterate, don't expect magic", "The real power is in the follow-up prompts — refine, drill down, reshape."),
        ("Tools are available today", "Everything we used is available right now — most with free tiers."),
        ("Start with one task", "Pick the most repetitive part of your week. Try AI on it tomorrow."),
    ]
    
    txb = tb(slide, Inches(1.1), Inches(1.6), Inches(11), Inches(5.5))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, (title, desc) in enumerate(lessons):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(16)
        p.text = f"▸  {title}"
        p.font.size = Pt(22)
        p.font.color.rgb = ORANGE
        p.font.bold = True
        p.font.name = FONT_NAME
        
        p2 = tf.add_paragraph()
        p2.text = f"    {desc}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = FONT_NAME
        p2.space_before = Pt(2)
    
    # ── Q&A Slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    txb = tb(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2))
    set_text(txb.text_frame, "Questions?", size=54, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.8), Inches(2.3), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()
    
    txb2 = tb(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(1.0))
    set_text(txb2.text_frame, "\"AI won't replace you. But someone using AI might.\nThat's why we're here tonight.\"", size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    txb2.text_frame.paragraphs[0].font.italic = True
    
    # ── Contact/Resources Slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.6), Inches(0.08), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()
    
    txb = tb(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.9))
    set_text(txb.text_frame, "Resources & Next Steps", size=36, bold=True, color=WHITE)
    
    resources = [
        ("Try Tonight (Free)", None),
        ("claude.ai", "Best for documents, analysis, and reasoning"),
        ("chatgpt.com", "Best all-rounder, image generation"),
        ("perplexity.ai", "Research with cited sources"),
        ("gamma.app", "Turn text into designed presentations"),
        ("", None),
        ("Get in Touch", None),
        ("Enriched Solutions", "enrichedsolutions.co.uk"),
    ]
    
    txb = tb(slide, Inches(1.1), Inches(1.6), Inches(10), Inches(5.5))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for name, desc in resources:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(8)
        
        if desc is None and name:
            # Section header
            p.text = name
            p.font.size = Pt(22)
            p.font.color.rgb = ORANGE
            p.font.bold = True
            p.font.name = FONT_NAME
            p.space_before = Pt(20)
        elif name == "":
            p.text = ""
            p.font.size = Pt(8)
        else:
            p.text = f"▸  {name}"
            p.font.size = Pt(20)
            p.font.color.rgb = CYAN
            p.font.bold = True
            p.font.name = FONT_NAME
            if desc:
                p2 = tf.add_paragraph()
                p2.text = f"    {desc}"
                p2.font.size = Pt(16)
                p2.font.color.rgb = LIGHT_GRAY
                p2.font.name = FONT_NAME
                p2.space_before = Pt(2)
    
    # Tagline at bottom
    txb = tb(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5))
    set_text(txb.text_frame, '"You cannot change the wind, but working together, we can adjust the sails"', size=14, color=RGBColor(0x88, 0x88, 0x99), alignment=PP_ALIGN.CENTER)
    txb.text_frame.paragraphs[0].font.italic = True
    
    return prs


# ── Main ──────────────────────────────────────────────────────────
if __name__ == "__main__":
    out_dir = os.path.dirname(os.path.abspath(__file__))
    
    print("Creating Talk 1: AI Overview...")
    prs1 = create_talk1()
    path1 = os.path.join(out_dir, "talk-1-ai-overview.pptx")
    prs1.save(path1)
    print(f"  ✓ Saved: {path1}")
    print(f"  Slides: {len(prs1.slides)}")
    
    print("\nCreating Talk 2: Workshop...")
    prs2 = create_talk2()
    path2 = os.path.join(out_dir, "talk-2-workshop.pptx")
    prs2.save(path2)
    print(f"  ✓ Saved: {path2}")
    print(f"  Slides: {len(prs2.slides)}")
    
    print("\n✅ Done!")
